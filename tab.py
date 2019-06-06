# -*- coding:utf-8 -*-

"""
@author: zhouxk
@time: 2019/6/3 17:26
@desc:
"""
from pptx import Presentation
from pptx.util import Inches

SLD_LAYOUT_TITLE_AND_CONTENT = 1
ppt = Presentation('xp.pptx')
first_page = ppt.slides[0]
title = '霞浦核电调试双周例会'
first_page.placeholders[0].text = title
date_time = '2019年06月10日'
first_page.placeholders[1].text = date_time

second_page = ppt.slides[1]
political_title = '信息主题：六年六组数据看全国查处违规公款吃喝'
second_page.placeholders[0].text = political_title
political_dept = '纪检监察处/审计处     李志平'
second_page.placeholders[1].text = political_dept
second_page.shapes[2].text = date_time

thrid_page = ppt.slides[2]
security_title = '信息主题：无XXXXX'
thrid_page.placeholders[0].text = political_title
security_dept = '生产准备处'
thrid_page.placeholders[1].text = security_dept
thrid_page.shapes[2].text = date_time

five_page = ppt.slides[4]
table = five_page.shapes[1]



tab_data = [{
    '序号': '1',
    '行动项': '同CIAE、CNPE和728院等几家设计院沟通，确定综合试验项目的必要性。',
    '责任部门': '运行处',
    '完成时间': '2019-04-30',
    '当前进展': '4月份分别和CIAE、CNPE两家设计院负责CFR600综合试验项目的负责人进行了交流、讨论。两家设计院的综合试验负责人从安全性、必要性两个方面对我们筛选的15个综合试验项目及试验平台提供了意见和建议。经王总同意，本行动项已延期至5月31日完成。'
}, {
    '序号': '2',
    '行动项': '各领域完成调试专用工器具清单和调试耗材清单梳理，清单提报生产准备处维修服务',
    '责任部门': '生产准备处',
    '完成时间': '2019-04-30',
    '当前进展': '陆续收到各部门提交数据，汇总后将审查意见反馈提报部门。后续将整理清单'
}, {
    '序号': '3',
    '行动项': '编制“钠品质控制专项”工作方案，实现对清洁度的全方位管控。',
    '责任部门': '调试管理处',
    '完成时间': '2019-05-31',
    '当前进展': '已根据各专业反馈意见完成方案修订'
}]

tab_frame = table.insert_table(rows=4, cols=5).table
tab_frame.columns[0].width = Inches(0.8)
tab_frame.columns[1].width = Inches(3.6)
tab_frame.columns[2].width = Inches(1.5)
tab_frame.columns[3].width = Inches(1.8)
tab_frame.columns[4].width = Inches(5)

first_title = []
content = []
for row in tab_data:
    arow = []
    for item in row.keys():
       first_title.append(item)
       arow.append(row[item])
    content.append(arow)

for i in range(4):
    for j in range(5):
        if i == 0 :
            tab_frame.cell(i, j).text = first_title[j]
        else:
            tab_frame.cell(i, j).text = content[i-1][j]

ppt.save('test.pptx')









# from pptx import Presentation
#
# prs = Presentation()
# bullet_slide_layout = prs.slide_layouts[1]
#
# slide = prs.slides.add_slide(bullet_slide_layout)
# shapes = slide.shapes
#
# title_shape = shapes.title
# body_shape = shapes.placeholders[1]
#
# title_shape.text = 'Adding a Bullet Slide'
#
# tf = body_shape.text_frame
# tf.text = 'Find the bullet slide layout Find the bullet slide layout '
#
# p = tf.add_paragraph()
# p.text = 'Use _TextFrame.text for first bullet'
# p.level = 1

# p = tf.add_paragraph()
# p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
# p.level = 2
#
# prs.save('test.pptx')













# from pptx import Presentation
#
# prs = Presentation()  # 创建对象
# title_slide_layout = prs.slide_layouts[0]   # 创建布局
# slide = prs.slides.add_slide(title_slide_layout)    # 用布局创建新页
# title = slide.shapes.title      # 新页形状
# subtitle = slide.placeholders[1]    # 新页占位符
#
# title.text = "Hello, World!"
# subtitle.text = "python-pptx was here!"
#
# prs.save('test.pptx')